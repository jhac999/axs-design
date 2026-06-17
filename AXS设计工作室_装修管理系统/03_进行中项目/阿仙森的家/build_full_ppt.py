import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import json

def create_full_ppt():
    prs = Presentation()
    # 强制 16:9 宽屏比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 10, 12)
    
    # Title Box
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "AXS 最终图文提案演示"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "阿仙森全案极客宣讲系统"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(212, 175, 55)

    # Read dynamic data from JSON
    dynamic_style = "基于深度心理数据定制的绝对安全空间 (3000K | 2.8m原顶 | 零人物)"
    json_path = "阿仙森审美与实用需求深度探测引擎收集表.json"
    if os.path.exists(json_path):
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                if "核心设计指令" in data and len(data["核心设计指令"]) > 0:
                    # 获取第一条指令作为定调
                    dynamic_style = data["核心设计指令"][0]
        except Exception as e:
            pass

    p3 = tf.add_paragraph()
    p3.text = f"\nAXS 提案信条：全程以硬核机电数据、环境心理学与大师理论进行空间推演。\n全局定调：{dynamic_style}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(150, 150, 150)

    slides_data = [
        {
            "img": "00_整体空间鸟瞰图.png",
            "tag": "Part 0 | 全局动线透视",
            "title": "上帝视角的空间骨架",
            "desc": "“在进入具体空间前，我们先用上帝视角俯瞰整个建筑的平面排布。这是所有的墙体拆改与洄游动线的根基。”"
        },
        {
            "img": "01_客厅全景_效果图.png",
            "tag": "Part 1 | 核心主线 LDK",
            "title": "消灭空旷，降维造茧",
            "desc": "“请特别注意顶面尺度。我们下了死命令：保留 2.8m 原顶找平，放弃所有暗藏筒灯，全部改用极简明装轨道灯矩阵。全案照明色温死死锁定在 3000K 均质暖白光。”"
        },
        {
            "img": "02_客餐厅一体_效果图.png",
            "tag": "Part 1 | 洞察与共鸣",
            "title": "为“茧居者”打造庇护所",
            "desc": "“基于您的物理数据模型，这套方案我们只做了一件事：做减法。我们剥离了 90% 的无效噪音，为您构建了一个绝对安全的物理系统。”"
        },
        {
            "img": "03_入户玄关_效果图.png",
            "tag": "Part 2 | 核心主线 归家",
            "title": "零内耗动线",
            "desc": "“强行把‘脱鞋-喝水-躺下’压缩在两步之内，引入直饮水路。3cm的微下沉落尘区，致敬本间贵史的安全理念，做足生命周期防线。”"
        },
        {
            "img": "04_主卧_效果图.png",
            "tag": "Part 2 | 核心主线 休憩",
            "title": "绝对静音舱",
            "desc": "“在这个空间里，没有任何多余的线条，只有纯粹的原木轻奢氛围与绝对深度的睡眠。”"
        },
        {
            "img": "05_厨房_效果图.png",
            "tag": "Part 2 | 核心主线 烹饪",
            "title": "死磕人体工学",
            "desc": "“厨房操作台做了 10 公分的高低台落差，水槽高、灶台低。做饭不再是纯粹的体力活。”"
        },
        {
            "img": "06_公卫_效果图.png",
            "tag": "Part 2 | 核心主线 卫浴",
            "title": "极致干湿分离",
            "desc": "“公卫完全执行干湿彻底分离，确保最高频洗漱场景的动线效率最大化。”"
        },
        {
            "img": "07_次卧_效果图.png",
            "tag": "Part 3 | 选配空间矩阵",
            "title": "生命周期可变空间",
            "desc": "“用全生命周期设计的眼光看待这间房，柜体全部采用 ENF 级零甲醛添加的顶级板材，守护最底层的健康红线。”"
        },
        {
            "img": "08_主卫_效果图.png",
            "tag": "Part 3 | 选配空间矩阵",
            "title": "Spa 级微水泥避难所",
            "desc": "“主卧套房内的绝对私密领地。大面积微水泥无缝包裹，壁挂马桶隐形水箱，消灭了一切视觉杂乱与卫生死角。”"
        },
        {
            "img": "09_阳台_效果图.png",
            "tag": "Part 3 | 选配空间矩阵",
            "title": "复合休闲家政区",
            "desc": "“洗烘套装完美隐藏，同时保留观景功能，将原本割裂的阳台彻底融入室内。”"
        },
        {
            "img": "10_书房_效果图.png",
            "tag": "Part 3 | 选配空间矩阵",
            "title": "声学阻断书房",
            "desc": "“具备极高的声学阻断力。纯实木悬浮长桌，配合全遮光厚重窗帘，展现了完美的向内收敛性。”"
        },
        {
            "img": "11_衣帽间_效果图.png",
            "tag": "Part 3 | 选配空间矩阵",
            "title": "高频收纳矩阵",
            "desc": "“精细化的收纳矩阵，防尘灯光系统，让每一次挑选衣物都充满仪式感。”"
        },
        {
            "img": "12_电视特写_效果图.png",
            "tag": "Part 4 | 极致工艺交底",
            "title": "毫米级材质碰撞",
            "desc": "“原石与木饰面的无缝拼缝，暗藏洗墙灯带毫米级退晕。这不仅是美学，更是我们在工地极度严苛的施工准则。”"
        },
        {
            "img": "13_床头特写_效果图.png",
            "tag": "Part 4 | 极致工艺交底",
            "title": "床头光影退晕",
            "desc": "“感受主卧核心背景的氛围光影退晕、软包与木饰面毫米级的拼接逻辑。”"
        },
        {
            "img": "14_餐边柜特写_效果图.png",
            "tag": "Part 4 | 极致工艺交底",
            "title": "水吧极限吞吐量",
            "desc": "“验证小家电操作区高度及收纳极限吞吐量设计，暗埋管线做到零外露。”"
        },
        {
            "img": "15_720漫游_效果图.png",
            "tag": "Part 4 | 满配沉浸漫游",
            "title": "720° VR 空间体验",
            "desc": "“这是最后的 VR 漫游。所有的唯美视觉，都严格建立在 AutoCAD 毫米级坐标工程图上。我们用绝对客观的数据，为您落地最主观的感官。”"
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 10, 12)
        
        img_path = data["img"]
        # 左侧放置图像，高度约束，保持居中感
        if os.path.exists(img_path):
            try:
                # 给图片 8.5 英寸宽度
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(8.5))
            except:
                pass
                
        # 右侧放置文本
        txBox = slide.shapes.add_textbox(Inches(9.4), Inches(2.0), Inches(3.5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = data["tag"]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(180, 180, 180)
        
        p2 = tf.add_paragraph()
        p2.text = data["title"]
        p2.font.bold = True
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(212, 175, 55)
        
        p3 = tf.add_paragraph()
        p3.text = "\n" + data["desc"]
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(230, 230, 230)

    output_path = '阿仙森_满配版极客宣讲_16张大满贯.pptx'
    prs.save(output_path)
    print(f"[OK] {output_path} Generated successfully.")

if __name__ == '__main__':
    create_full_ppt()
