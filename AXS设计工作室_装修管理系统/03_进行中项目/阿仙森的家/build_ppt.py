from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

BASE_DIR = r"f:\吉胡阿川\01lhjk\事业\AXS设计工作室\AXS设计工作室_装修管理系统\03_进行中项目\阿仙森的家"

def create_ppt():
    print("[SYSTEM] Booting AXS Commercial PPTX Engine (Master Cited)...")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_layout = prs.slide_layouts[6]
    
    def set_black_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(12, 12, 12)

    # ================= PART 1: 洞察与共鸣 (Insight) =================
    slide_1 = prs.slides.add_slide(blank_layout)
    set_black_bg(slide_1)
    
    txBox = slide_1.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "为“茧居型松弛者”打造绝对宁静的庇护所"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "\n[痛点解读] 极度抗拒强光刺激、归家能量极速衰减\n[主调包装] 极简有机 / 原木轻奢 (Japandi)"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(180, 180, 180)
    
    p3 = tf.add_paragraph()
    p3.text = "\n[审美与机电红线声明]\n1. 坚决舍弃压迫层高的平吊顶，全案2.8m原顶找平，明装极简灯具。\n2. 【严防代入感偏差】所有空间画面严禁出现人物，保障纯粹的空间尺度。\n3. 全局光照锚定在 3000K 均质暖白，剔除一切光斑割裂感。"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(255, 80, 80)

    # ================= PART 2: 核心痛点靶向解法 =================
    target_slides = [
        {
            "img": "客餐厅_明装射灯版效果图.png",
            "psy_title": "场景A：消灭空旷，降维造茧 (客厅核心区)",
            "psy_desc": "彻底放弃大横厅带来的视觉压力。通过极简明装轨道灯与原木轻奢材质包裹，复刻小房间的安全感，支持进门后极致放空半躺。",
            "sys_title": "[工程数据面板] 空间与机电参数",
            "sys_desc": "净空保留：2800mm (拒绝平吊顶)\n主光源色温：全局 3000K 暖白漫反射\n隐藏收纳体积：> 12m³ (墙面 0 杂物)"
        },
        {
            "img": "02_客餐厅一体_效果图.png",
            "psy_title": "场景A-2：洄游动线与全局视觉控制 (LDK)",
            "psy_desc": "超大广角展现客餐厅的无缝穿透。严禁出现任何多余的实体隔断，纯靠原木材质的延伸与天花板的留白来划分区域边界，维持绝对纯净的建筑尺度。",
            "sys_title": "[工程数据面板] 动线与光影参数",
            "sys_desc": "视觉穿透率：100% (零视线阻挡)\n顶面平整度：原生楼板直刷无机涂料\n全局照度均匀度：U0 > 0.8"
        },
        {
            "img": "入户玄关_明装射灯版效果图.png",
            "psy_title": "场景B：零内耗动线与能量回血",
            "psy_desc": "深刻共情您的能量衰减模型：“脱鞋 -> 喝水 -> 躺下”被压缩进仅有两步的极限距离内，切断传统走向厨房倒水的漫长动线。",
            "sys_title": "[工程数据面板] 拓扑参数",
            "sys_desc": "落尘区微降：-30mm\n机电改动：玄关特批引入直饮水路 (10L/h)\n动线物理缩短率：45%"
        }
    ]

    for index, data in enumerate(target_slides):
        slide = prs.slides.add_slide(blank_layout)
        set_black_bg(slide)
        
        img_path = os.path.join(BASE_DIR, data["img"])
        if os.path.exists(img_path):
            try:
                pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), height=Inches(9))
            except Exception:
                pass
                
        txBox_psy = slide.shapes.add_textbox(Inches(10), Inches(0.5), Inches(5.5), Inches(4))
        txBox_psy.text_frame.word_wrap = True
        p_psy_t = txBox_psy.text_frame.paragraphs[0]
        p_psy_t.text = data["psy_title"]
        p_psy_t.font.size = Pt(28)
        p_psy_t.font.bold = True
        p_psy_t.font.color.rgb = RGBColor(201, 168, 76)
        
        p_psy_d = txBox_psy.text_frame.add_paragraph()
        p_psy_d.text = "\n" + data["psy_desc"]
        p_psy_d.font.size = Pt(20)
        p_psy_d.font.color.rgb = RGBColor(220, 220, 220)

        txBox_sys = slide.shapes.add_textbox(Inches(10), Inches(4.8), Inches(5.5), Inches(4))
        txBox_sys.text_frame.word_wrap = True
        p_sys_t = txBox_sys.text_frame.paragraphs[0]
        p_sys_t.text = data["sys_title"]
        p_sys_t.font.size = Pt(24)
        p_sys_t.font.bold = True
        p_sys_t.font.color.rgb = RGBColor(255, 80, 80)
        
        p_sys_d = txBox_sys.text_frame.add_paragraph()
        p_sys_d.text = "\n" + data["sys_desc"]
        p_sys_d.font.size = Pt(18)
        p_sys_d.font.color.rgb = RGBColor(0, 255, 128) 

    # ================= PART 3: 隐形资产与大师溯源 =================
    slide_3 = prs.slides.add_slide(blank_layout)
    set_black_bg(slide_3)
    txBox3 = slide_3.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(7))
    tf3 = txBox3.text_frame
    tf3.paragraphs[0].text = "核心背书：立足于全球大师底层的【原木轻奢】体系"
    tf3.paragraphs[0].font.size = Pt(40)
    tf3.paragraphs[0].font.color.rgb = RGBColor(201, 168, 76)
    
    p_3_1 = tf3.add_paragraph()
    p_3_1.text = "\n【溯源认证 1】深泽直人 (Naoto Fukasawa) —— 《无意识设计 (Without Thought)》"
    p_3_1.font.size = Pt(24)
    p_3_1.font.bold = True
    p_3_1.font.color.rgb = RGBColor(255, 255, 255)
    
    p_3_1_desc = tf3.add_paragraph()
    p_3_1_desc.text = "证据溯源：大师在其著作中提出“好的设计是融入潜意识而消失的”。\n本案物理印证：我们将“脱鞋-喝水-躺下”极限压缩，玄关直接引入水吧，让喝水不经过任何大脑思考，极度符合原木轻奢的核心奥义。"
    p_3_1_desc.font.size = Pt(18)
    p_3_1_desc.font.color.rgb = RGBColor(200, 200, 200)

    p_3_2 = tf3.add_paragraph()
    p_3_2.text = "\n【溯源认证 2】本间贵史 (Takafumi Homma) —— 极致生命周期防线"
    p_3_2.font.size = Pt(24)
    p_3_2.font.bold = True
    p_3_2.font.color.rgb = RGBColor(255, 255, 255)

    p_3_2_desc = tf3.add_paragraph()
    p_3_2_desc.text = "证据溯源：在《梦想改造家》中多次强调“零高差与微死角的消除”。\n本案物理印证：入户区极限下沉控制在 3cm（防尘且防绊倒）；全案效果图[严禁人物出现]，剔除干扰，让极度纯净的空间直接对接受老龄与宠物的无障碍底线。"
    p_3_2_desc.font.size = Pt(18)
    p_3_2_desc.font.color.rgb = RGBColor(200, 200, 200)

    # ================= PART 5: 商业底牌与落地清单 =================
    slide_5 = prs.slides.add_slide(blank_layout)
    set_black_bg(slide_5)
    
    txBox5 = slide_5.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(6))
    tf5 = txBox5.text_frame
    tf5.paragraphs[0].text = "商业与系统底牌"
    tf5.paragraphs[0].font.size = Pt(44)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    p_5_1 = tf5.add_paragraph()
    p_5_1.text = "\n【出图矩阵清单交付】\n标配 6 大核心区（客厅/玄关/主卧等）已锁定方案骨架。\n（次卧、独立衣帽间及全屋720漫游按标准化增项计费）"
    p_5_1.font.size = Pt(24)
    p_5_1.font.color.rgb = RGBColor(200, 200, 200)

    p_5_2 = tf5.add_paragraph()
    p_5_2.text = "\n【全案落地商业承诺】\n基装材料、软装家具，AXS 系统供应链底价 100% 开放。\n没有材料差价，没有暗箱增项。全盘仅收合同金额 10% 极客纯服务费。"
    p_5_2.font.size = Pt(32)
    p_5_2.font.bold = True
    p_5_2.font.color.rgb = RGBColor(201, 168, 76)
    
    output_path = os.path.join(BASE_DIR, "阿仙森_原木轻奢版商业汇报.pptx")
    prs.save(output_path)
    print(f"[OK] Commercial PPTX Engine Executed Successfully.")

if __name__ == "__main__":
    create_ppt()
