from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

BASE_DIR = r"F:\吉胡阿川\01lhjk\事业\AXS设计工作室\AXS设计工作室_装修管理系统\03_进行中项目\格哥的空间\20260531_02_工程骨架与概念出图"

def create_ppt():
    print("[SYSTEM] Booting AXS Triple-Model PPTX Engine (11-Image Full Roster)...")
    prs = Presentation()
    
    # 强制 16:9 比例 (宽屏)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    def set_black_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(12, 12, 12) # 深邃极黑

    # ================= 幻灯片 1：封面 =================
    slide_1 = prs.slides.add_slide(blank_slide_layout)
    set_black_bg(slide_1)
    
    txBox = slide_1.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AXS 终极系统诊断与物理推演"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "\n基于《17题潜意识扫描与物理机电压测》的靶向输出报告\n代号：格哥的宁静庇护所"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(201, 168, 76) # AXS 金色

    # ================= 三重模型数据 (11 张满编制) =================
    slides_data = [
        # 客厅 (3张)
        {
            "img": "gege_render_living_tv_1780217831312.png",
            "psy_title": "■ 心理抚慰：去装饰化留白",
            "psy_desc": "您在问卷中提到：'最怕下班回家看到乱七八糟的杂物'。\n\n我们摒弃了一切繁琐的电视背景墙与硬包，采用全屋微水泥实现绝对的视觉统一。用极致的空旷，来承载您对“宁静”的渴望。",
            "sys_title": "▶ 建筑红线与机电参数",
            "sys_desc": "[空间约束] 顶面零跌级大平顶，净空强保 2850mm\n[光照锁死] 取消主灯，防眩光深杯射灯，色温 3500K\n[材质系统] 墙地一体化微水泥涂层，去踢脚线工艺\n[环控预留] 隐形风口，风管机下出下回循环系统"
        },
        {
            "img": "gege_render_living_sofa_1780217844777.png",
            "psy_title": "■ 心理抚慰：消灭视觉噪音",
            "psy_desc": "针对您【视觉洁癖容忍度极低】的潜意识特征。\n\n我们将所有可能引发焦虑的日常杂物（纸巾、钥匙、快递箱）全部抹除。让沙发背景永远保持空旷的通透感。",
            "sys_title": "▶ 容积压测与吞吐参数",
            "sys_desc": "[隐形容积] 强制切出隐藏收纳矩阵\n[通透比例] LDK 客餐厅视线无阻碍穿透 > 8.5m\n[家具尺度] 超低矮模块化沙发，降低视觉重心 30%"
        },
        {
            "img": "gege_render_living_detail_1780234122312.png",
            "psy_title": "■ 心理抚慰：触手可及的温度",
            "psy_desc": "在冰冷极简的微水泥基底上，我们为您保留了最温润的交互界面。\n\n当您坐在沙发上，手边就是高级木饰面与微水泥的极致材料碰撞，用最细微的质感抚平外界带来的所有粗糙。",
            "sys_title": "▶ 材质级细节与光影红线",
            "sys_desc": "[光影解算] 3000K 洗墙灯槽，退晕光斑完美控制在 50cm 内\n[材质碰撞] 西班牙哑光微水泥 vs 科技木饰面\n[折射率] 所有表面漫反射率 < 10%，彻底杜绝刺眼反光"
        },
        # 餐厅与玄关 (2张)
        {
            "img": "gege_render_dining_1780217857647.png",
            "psy_title": "■ 心理抚慰：无压力简餐区",
            "psy_desc": "您提到需要一个随时能喝杯咖啡、吃简餐的无压迫角落。\n\n岛台与餐桌的几何咬合是最佳解药，在这里，一切动作都可以极度慵懒。",
            "sys_title": "▶ 动线拓扑与空间布局",
            "sys_desc": "[动线] 洄游动线宽度预留 900mm，保证绝对顺畅\n[容积] 背后满墙柜体提供 >1.8 m³ 隐性收纳\n[机电] 预埋 10A 轨道插座，满足小家电灵活供电"
        },
        {
            "img": "gege_render_entryway_1780234133566.png",
            "psy_title": "■ 心理抚慰：卸下铠甲的缓冲区",
            "psy_desc": "每天推开家门的第一眼，决定了您归家时的心率。\n\n我们将入户玄关打造成了一条极简隧道，光线指引您步入明亮的客厅。在换鞋的那 10 秒钟，您可以彻底卸下外面的伪装与疲惫。",
            "sys_title": "▶ 玄关吞吐量与视觉焦点",
            "sys_desc": "[悬浮工艺] 底部悬空 150mm，内藏感应灯带\n[吞吐量] 侧边隐藏式鞋柜支持收纳 45 双标准尺寸鞋\n[视觉延伸] 微水泥地面无缝延伸至全屋，抹除空间割裂感"
        },
        # 主卧套房 (3张)
        {
            "img": "gege_render_master_bed_1780217880559.png",
            "psy_title": "■ 心理抚慰：深度睡眠舱",
            "psy_desc": "您说“关上门谁都别找我，想一个人待一会儿”。\n\n我们在主卧套房内，利用深木饰面的排版和单独的光影控制，为您强行扣出了一个独立避难所。当您躺在这里，门外的喧闹将与您彻底隔绝。",
            "sys_title": "▶ 声学隔离与睡眠系统",
            "sys_desc": "[声学红线] 主卧隔墙强制加厚 50mm\n[材质阻断] 墙体内嵌高密度吸音棉，隔音标准 STC > 50\n[智控预设] 一键切断主光源，仅保留 10lx 起夜地脚灯"
        },
        {
            "img": "gege_render_master_walkway_1780234149247.png",
            "psy_title": "■ 心理抚慰：丝滑的晨间仪式",
            "psy_desc": "早晨起床、更衣、洗漱，这段时间内不能有任何物理阻碍。\n\n我们打通了衣帽间与主卫的边界，取消了所有的传统房门，让您的每一个晨间动作都如流水般顺畅且高效。",
            "sys_title": "▶ 步态分析与气口调度",
            "sys_desc": "[无障碍通行] 睡眠区至卫浴间 0mm 高差\n[藏风聚气] 床头避开气口直冲，符合环境心理学\n[系统收纳] 悬浮衣帽间提供 >3.5 m³ 的绝对挂衣体量"
        },
        {
            "img": "gege_render_master_corner_1780234161178.png",
            "psy_title": "■ 心理抚慰：2㎡ 精神角落",
            "psy_desc": "哪怕是在主卧这个绝对私密的空间里，您依然需要一个甚至连伴侣都无法打扰的角落。\n\n一把舒适的单椅，一盏极简的落地灯，这就是您每晚睡前抽根烟、喝杯酒、彻底放空大脑的精神孤岛。",
            "sys_title": "▶ 人体工学与照度模拟",
            "sys_desc": "[光域限制] 落地灯照度被死死锁在半径 1m 内，不干扰伴侣\n[视角拦截] 靠窗视角，切断与室内主通道的视觉交叉\n[地面温控] 此区域铺设强化发热电缆，保障赤足舒适度"
        },
        # 厨房 (1张)
        {
            "img": "gege_render_kitchen_1780217892457.png",
            "psy_title": "■ 心理抚慰：家务隔离区",
            "psy_desc": "针对您排斥家务的底层诉求。\n\n我们将油污区与清洗区进行硬性切分，打造最高效的洗切炒流水线。在这里，烹饪不再是苦役，而是掌控一切的实验。",
            "sys_title": "▶ 流水线排布与重机电矩阵",
            "sys_desc": "[操作红线] 操作台总长度 > 3.2m\n[排烟系统] 22m³/min 爆炒级风量预留，顶装止逆阀\n[水电核载] 专线 4平方线缆直达，支撑烤箱与洗碗机同时开启"
        },
        # 次卧 (1张)
        {
            "img": "gege_render_guest_bed_1780234174957.png",
            "psy_title": "■ 心理抚慰：不被冷落的次级空间",
            "psy_desc": "即便是给父母或客人偶尔留宿的房间，也绝不能与整体系统产生割裂。\n\n我们将整体的高级灰基调与极简骨架延续至此，用最低的成本，实现了最高级的空间质感统一。",
            "sys_title": "▶ 造价控制与标准化落位",
            "sys_desc": "[成本熔断] 砍掉一切非必要吊顶，仅保留窗帘盒与边板\n[整合设计] 床、衣柜、书桌 100% 模数化统一定制\n[温控标准] 同样接入全屋冷暖环控系统，无体验断层"
        },
        # 卫生间 (1张)
        {
            "img": "gege_render_bathroom_1780217904892.png",
            "psy_title": "■ 心理抚慰：情绪释放带",
            "psy_desc": "在最私密、最容易产生水渍焦虑的空间。\n\n我们用整块大尺度镜柜与悬浮台面，消灭了所有的卫生死角。洗漱时的每一滴水，都在我们的系统控制之内。",
            "sys_title": "▶ 防死角工艺与防潮防线",
            "sys_desc": "[防水底线] 强制要求满刷至顶 1.8m\n[工艺红线] 悬浮台下盆无缝拼接，墙地砖十字对缝误差 <1mm\n[收容矩阵] 隐形大尺寸镜柜，吞吐量 >0.5 m³"
        }
    ]
    
    # 循环生成【三重模型】幻灯片
    for index, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        set_black_bg(slide)
        
        # 0. 幻灯片编号水印
        txBox_num = slide.shapes.add_textbox(Inches(13.5), Inches(0.2), Inches(2), Inches(1))
        tf_num = txBox_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = f"NO.{index+1:02d}"
        p_num.font.size = Pt(40)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(100, 100, 100) # 暗灰色水印
        
        # 1. 左侧：美学图纸 (60% 宽度，即 9.6 英寸)
        img_path = os.path.join(BASE_DIR, data["img"])
        if os.path.exists(img_path):
            try:
                pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0), height=Inches(9))
            except Exception as e:
                print(f"[Error] Image insertion failed {img_path}: {e}")
        else:
            print(f"[Warning] Image not found: {img_path}")
            
        # 2. 右上：心理文案 (左边距 10, 顶部 0.5, 宽 5.5, 高 4)
        txBox_psy = slide.shapes.add_textbox(Inches(10), Inches(0.5), Inches(5.5), Inches(4))
        txBox_psy.text_frame.word_wrap = True
        
        p_psy_t = txBox_psy.text_frame.paragraphs[0]
        p_psy_t.text = f"NO.{index+1:02d} {data['psy_title']}"
        p_psy_t.font.size = Pt(24)
        p_psy_t.font.bold = True
        p_psy_t.font.color.rgb = RGBColor(201, 168, 76) # 金色
        
        txBox_psy.text_frame.add_paragraph()
        p_psy_d = txBox_psy.text_frame.add_paragraph()
        p_psy_d.text = data["psy_desc"]
        p_psy_d.font.size = Pt(18)
        p_psy_d.font.color.rgb = RGBColor(220, 220, 220)

        # 3. 右下：硬核数据 (左边距 10, 顶部 4.8, 宽 5.5, 高 4)
        txBox_sys = slide.shapes.add_textbox(Inches(10), Inches(4.8), Inches(5.5), Inches(4))
        txBox_sys.text_frame.word_wrap = True
        
        p_sys_t = txBox_sys.text_frame.paragraphs[0]
        p_sys_t.text = data["sys_title"]
        p_sys_t.font.size = Pt(22)
        p_sys_t.font.bold = True
        p_sys_t.font.color.rgb = RGBColor(255, 80, 80) # 红色警示
        
        txBox_sys.text_frame.add_paragraph()
        p_sys_d = txBox_sys.text_frame.add_paragraph()
        p_sys_d.text = data["sys_desc"]
        p_sys_d.font.size = Pt(16)
        p_sys_d.font.color.rgb = RGBColor(0, 255, 128) 

    # ================= 幻灯片尾页：商业底线 =================
    slide_end = prs.slides.add_slide(blank_slide_layout)
    set_black_bg(slide_end)
    
    txBox = slide_end.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(4))
    tf = txBox.text_frame
    
    p = tf.paragraphs[0]
    p.text = "AXS 商业防御底线"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = tf.add_paragraph()
    p2.text = "\n本案仅收合同 10% 的极客纯服务费。\n基装材料、软装家具，AXS 系统供应链底价 100% 开放给主理人。"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(180, 180, 180)
    
    p3 = tf.add_paragraph()
    p3.text = "\n没有材料差价，没有增项回扣。用绝对的透明，换取绝对的信任。"
    p3.font.size = Pt(36)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(201, 168, 76) 
    
    # 保存文件
    output_path = os.path.join(BASE_DIR, "AXS001 格哥的家_满编制版.pptx")
    prs.save(output_path)
    print(f"[OK] Triple-Model PPTX Engine Executed Successfully.")
    print(f"[Output] -> {output_path}")

if __name__ == "__main__":
    create_ppt()
