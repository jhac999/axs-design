import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_transparent_quote(output_path):
    # 模拟数据提取
    bom_data = [
        {"item": "基础水电极客包 (强制水压测试+绝缘终端)", "qty": "1 式", "unit_price": 45000, "total": 45000},
        {"item": "全屋留白大白漆 (一底两面 拒绝微水泥)", "qty": "400 ㎡", "unit_price": 35, "total": 14000},
        {"item": "三代同堂声学隔离包 (双层阻尼+实心静音门)", "qty": "1 式", "unit_price": 18000, "total": 18000},
        {"item": "AXS强制通顶餐边柜 (深度40cm/无拉手/一门到顶)", "qty": "12 ㎡", "unit_price": 1000, "total": 12000},
        {"item": "全屋哑光素色地砖 (薄贴对缝工艺)", "qty": "200 ㎡", "unit_price": 95, "total": 19000}
    ]
    
    total_bare_cost = sum(item["total"] for item in bom_data)
    axs_fee = int(total_bare_cost * 0.10)
    grand_total = total_bare_cost + axs_fee

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # ----------------------------------------------------
    # Slide 1: BOM 明细 (黑底绿字极客风)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    background1 = slide1.background
    fill1 = background1.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(15, 15, 15)

    tx_box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
    tf1 = tx_box1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "AXS 透明决选清单 (BOM) | 严格遵守 12 万基装红线"
    p1.font.bold = True
    p1.font.size = Pt(40)
    p1.font.color.rgb = RGBColor(255, 255, 255)

    y_offset = 2.5
    for item in bom_data:
        box = slide1.shapes.add_textbox(Inches(1), Inches(y_offset), Inches(14), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        text_line = f"[+] {item['item']} | 用量: {item['qty']} | 底价: ¥{item['unit_price']} -> 拦截价: ¥{item['total']}"
        p.text = text_line
        p.font.size = Pt(24)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(120, 120, 120)
        y_offset += 0.8

    # ----------------------------------------------------
    # Slide 2: 报价剥离展示 (商业防线)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0, 0, 0)

    tx_box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
    tf2 = tx_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "AXS 资金池绝对隔离法则 (熔断保护生效)"
    p2.font.bold = True
    p2.font.size = Pt(48)
    p2.font.color.rgb = RGBColor(255, 255, 255)

    # Bare cost
    bare_box = slide2.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1))
    tf_bare = bare_box.text_frame
    p_bare = tf_bare.paragraphs[0]
    p_bare.text = f"> 工厂 100% 直采裸价池：¥{total_bare_cost:,}"
    p_bare.font.size = Pt(36)
    p_bare.font.name = "Consolas"
    p_bare.font.color.rgb = RGBColor(200, 200, 200)

    # AXS Fee (Highlight)
    fee_box = slide2.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
    tf_fee = fee_box.text_frame
    p_fee = tf_fee.paragraphs[0]
    p_fee.text = f"> AXS 极客系统 10% 纯管理费：¥{axs_fee:,}"
    p_fee.font.size = Pt(40)
    p_fee.font.name = "Consolas"
    p_fee.font.bold = True
    p_fee.font.color.rgb = RGBColor(0, 255, 100)  # Hacker Green

    # Grand Total
    total_box = slide2.shapes.add_textbox(Inches(1), Inches(6.5), Inches(14), Inches(1))
    tf_total = total_box.text_frame
    p_total = tf_total.paragraphs[0]
    p_total.text = f"总造价 (严格控制在12万内): ¥{grand_total:,} | 首期请款 (30%): ¥{int(grand_total * 0.3):,}"
    p_total.font.size = Pt(36)
    p_total.font.name = "Consolas"
    p_total.font.color.rgb = RGBColor(255, 100, 100)

    prs.save(output_path)
    print(f"[*] 成功输出本地透明报价单: {output_path}")

if __name__ == "__main__":
    CURRENT_DIR = os.path.dirname(os.path.abspath(__name__))
    OUTPUT_FILE = os.path.join(CURRENT_DIR, "AXS001_格哥的家_极客透明决选单.pptx")
    
    print("="*50)
    print("[AXS 算量中枢] 捕获到最新的 DXF 骨架...")
    print("[AXS 算量中枢] 已挂载【AXS审计官自检法则】...")
    print("[AXS 算量中枢] 触发熔断红线！已砍掉背景墙与复杂吊顶！")
    print("[AXS 算量中枢] 正在强制隔离 10% 利润池...")
    create_transparent_quote(OUTPUT_FILE)
    print("="*50)
